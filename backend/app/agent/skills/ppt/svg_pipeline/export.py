"""Export SVG page files to an editable pptx presentation.

MIT inspiration: ppt-master ``skills/ppt-master/scripts/svg_to_pptx/``.
Minerva v1 uses ElementTree XML parsing for ``<text>`` and ``<rect>`` nodes so
core export works without optional ``ppt-svg`` dependencies (svglib/reportlab).
"""

from __future__ import annotations

import re
import xml.etree.ElementTree as ET
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt

from app.agent.skills.ppt.shared.transitions import apply_slide_transition

_SLIDE_WIDTH_IN = 10.0
_SLIDE_HEIGHT_IN = 5.625
_TRANSLATE_RE = re.compile(
    r"translate\s*\(\s*([-\d.]+)(?:[,\s]+([-\d.]+))?\s*\)",
    re.IGNORECASE,
)


class SvgExportError(Exception):
    """Structured SVG export failure."""

    def __init__(self, code: str, message: str) -> None:
        """Store error code and human-readable message."""

        super().__init__(message)
        self.code = code
        self.message = message


def _local_tag(tag: str) -> str:
    """Return the XML local name without namespace prefix."""

    if "}" in tag:
        return tag.rsplit("}", 1)[-1]
    return tag


def _parse_float(value: str | None, default: float = 0.0) -> float:
    """Parse a numeric SVG attribute value."""

    if value is None:
        return default
    try:
        return float(str(value).strip())
    except ValueError:
        return default


def _parse_viewbox(root: ET.Element) -> tuple[float, float, float, float]:
    """Read viewBox or fall back to width/height attributes."""

    viewbox = (root.get("viewBox") or "").strip()
    if viewbox:
        parts = [p for p in viewbox.replace(",", " ").split() if p]
        if len(parts) == 4:
            return tuple(float(p) for p in parts)  # type: ignore[return-value]
    width = _parse_float(root.get("width"), 960.0)
    height = _parse_float(root.get("height"), 540.0)
    return 0.0, 0.0, width, height


def _parse_translate(transform: str | None) -> tuple[float, float]:
    """Extract translate(x, y) offsets from a transform attribute."""

    if not transform:
        return 0.0, 0.0
    match = _TRANSLATE_RE.search(transform)
    if not match:
        return 0.0, 0.0
    x = float(match.group(1))
    y = float(match.group(2) or 0.0)
    return x, y


def _element_position(element: ET.Element) -> tuple[float, float]:
    """Resolve x/y for an SVG element including simple translate transforms."""

    x = _parse_float(element.get("x"))
    y = _parse_float(element.get("y"))
    dx, dy = _parse_translate(element.get("transform"))
    return x + dx, y + dy


def _svg_units_to_inches(value: float, svg_extent: float, slide_extent_in: float) -> float:
    """Map SVG user units to slide inches."""

    if svg_extent <= 0:
        return 0.0
    return (value / svg_extent) * slide_extent_in


def _parse_hex_color(color: str | None, default: RGBColor | None = None) -> RGBColor | None:
    """Parse a #RRGGBB fill or stroke color."""

    if not color:
        return default
    raw = color.strip()
    if raw.lower() in {"none", "transparent"}:
        return default
    if raw.startswith("#"):
        raw = raw[1:]
    if len(raw) == 3:
        raw = "".join(ch * 2 for ch in raw)
    if len(raw) != 6:
        return default
    try:
        return RGBColor(int(raw[0:2], 16), int(raw[2:4], 16), int(raw[4:6], 16))
    except ValueError:
        return default


def _element_text(element: ET.Element) -> str:
    """Collect visible text from a text element and nested tspans."""

    parts: list[str] = []
    if element.text:
        parts.append(element.text.strip())
    for child in element:
        if _local_tag(child.tag) == "tspan" and child.text:
            parts.append(child.text.strip())
        if child.tail:
            parts.append(child.tail.strip())
    return " ".join(part for part in parts if part)


def _blank_slide_layout(prs: Presentation):
    """Pick a slide layout without placeholders."""

    for layout in prs.slide_layouts:
        if len(layout.placeholders) == 0:
            return layout
    return prs.slide_layouts[-1]


def _add_rect_shape(slide, element: ET.Element, vb_w: float, vb_h: float) -> None:
    """Add a rectangle autoshape from an SVG rect element."""

    x, y = _element_position(element)
    width = _parse_float(element.get("width"))
    height = _parse_float(element.get("height"))
    if width <= 0 or height <= 0:
        return

    left = Inches(_svg_units_to_inches(x, vb_w, _SLIDE_WIDTH_IN))
    top = Inches(_svg_units_to_inches(y, vb_h, _SLIDE_HEIGHT_IN))
    shape_width = Inches(_svg_units_to_inches(width, vb_w, _SLIDE_WIDTH_IN))
    shape_height = Inches(_svg_units_to_inches(height, vb_h, _SLIDE_HEIGHT_IN))

    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, shape_width, shape_height)
    shape.line.fill.background()
    color = _parse_hex_color(element.get("fill"))
    if color is not None:
        shape.fill.solid()
        shape.fill.fore_color.rgb = color


def _add_text_shape(slide, element: ET.Element, vb_w: float, vb_h: float) -> None:
    """Add an editable text box from an SVG text element."""

    text = _element_text(element)
    if not text:
        return

    x, y = _element_position(element)
    font_size = _parse_float(element.get("font-size") or element.get("fontSize"), 24.0)
    # SVG y is baseline; approximate textbox top from font size.
    top_y = max(y - font_size, 0.0)
    est_width = max(len(text) * font_size * 0.55, font_size * 2)
    est_height = max(font_size * 1.4, font_size)

    left = Inches(_svg_units_to_inches(x, vb_w, _SLIDE_WIDTH_IN))
    top = Inches(_svg_units_to_inches(top_y, vb_h, _SLIDE_HEIGHT_IN))
    box_width = Inches(_svg_units_to_inches(est_width, vb_w, _SLIDE_WIDTH_IN))
    box_height = Inches(_svg_units_to_inches(est_height, vb_h, _SLIDE_HEIGHT_IN))

    textbox = slide.shapes.add_textbox(left, top, box_width, box_height)
    text_frame = textbox.text_frame
    text_frame.clear()
    paragraph = text_frame.paragraphs[0]
    run = paragraph.add_run()
    run.text = text
    run.font.size = Pt(font_size * 0.75)
    color = _parse_hex_color(element.get("fill"))
    if color is not None:
        run.font.color.rgb = color


def _export_single_svg(prs: Presentation, svg_path: Path, *, transition: str) -> None:
    """Append one slide built from rects and text nodes in an SVG file."""

    try:
        tree = ET.parse(svg_path)
    except (ET.ParseError, OSError) as exc:
        raise SvgExportError("svg_parse_failed", f"failed to parse {svg_path.name}: {exc}") from exc

    root = tree.getroot()
    _, _, vb_w, vb_h = _parse_viewbox(root)
    if vb_w <= 0 or vb_h <= 0:
        raise SvgExportError("svg_parse_failed", f"invalid viewBox in {svg_path.name}")

    slide = prs.slides.add_slide(_blank_slide_layout(prs))

    for element in root.iter():
        tag = _local_tag(element.tag)
        if tag == "rect":
            _add_rect_shape(slide, element, vb_w, vb_h)
        elif tag == "text":
            _add_text_shape(slide, element, vb_w, vb_h)

    if transition != "keep":
        apply_slide_transition(slide, transition)


def export_svgs_to_pptx(
    svg_paths: list[Path],
    output_pptx: Path,
    *,
    transition: str = "fade",
) -> int:
    """Export SVG pages to editable pptx; returns slide count."""

    if not svg_paths:
        raise SvgExportError("svg_export_failed", "no SVG pages to export")

    prs = Presentation()
    prs.slide_width = Inches(_SLIDE_WIDTH_IN)
    prs.slide_height = Inches(_SLIDE_HEIGHT_IN)

    for svg_path in svg_paths:
        if not svg_path.is_file():
            raise SvgExportError("svg_export_failed", f"SVG page not found: {svg_path}")
        _export_single_svg(prs, svg_path, transition=transition)

    output_pptx.parent.mkdir(parents=True, exist_ok=True)
    try:
        prs.save(str(output_pptx))
    except OSError as exc:
        raise SvgExportError("write_failed", str(exc)) from exc

    return len(svg_paths)
