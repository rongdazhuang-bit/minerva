"""Orchestrate PPT generation from outline to output file."""

from __future__ import annotations

import json
import uuid
from pathlib import Path
from typing import Any

from langchain_core.language_models.chat_models import BaseChatModel
from pptx import Presentation

from app.agent.infrastructure.agent_file_sandbox import AgentFileSandbox
from app.agent.skills.ppt.pptmaker.constants import DEFAULT_LAYOUT_INDEX_PATH, DEFAULT_TEMPLATE_PATH
from app.agent.skills.ppt.pptmaker.fill import (
    build_placeholder_labels,
    enrich_labels_from_template,
    fill_slide_content,
    remove_all_existing_slides,
)
from app.agent.skills.ppt.pptmaker.layout_select import select_layout_name
from app.agent.skills.ppt.pptmaker.normalize import expand_outline_with_meta
from app.agent.skills.ppt.shared.notes import set_speaker_notes
from app.agent.skills.ppt.shared.transitions import apply_slide_transition
from app.agent.skills.ppt.template_fill.analyze import analyze_template
from app.agent.skills.ppt.template_fill.apply import apply_fill_plan
from app.agent.skills.ppt.template_fill.check_plan import check_fill_plan
from app.agent.skills.ppt.template_fill.plan_builder import outline_to_fill_plan

_IMAGE_SUFFIXES = {".png", ".jpg", ".jpeg", ".webp", ".gif", ".bmp"}


class PptGenerateError(Exception):
    """Structured generation failure."""

    def __init__(self, code: str, message: str) -> None:
        """Store error code and message."""

        super().__init__(message)
        self.code = code
        self.message = message


def _load_layout_index(path: Path) -> list[dict[str, Any]]:
    """Load layout index JSON from disk."""

    if not path.is_file():
        raise PptGenerateError("template_missing", f"layout index not found: {path}")
    raw = json.loads(path.read_text(encoding="utf-8"))
    if isinstance(raw, list):
        return raw
    if isinstance(raw, dict) and isinstance(raw.get("layouts"), list):
        return raw["layouts"]
    raise PptGenerateError(
        "template_missing",
        f"layout index has invalid format (expected a list): {path}",
    )


def _build_placeholder_geometry(layout_entry: dict[str, Any] | None) -> dict[str, dict[str, float]]:
    """Map placeholder label to geometry hints from a layout_index entry."""

    if not layout_entry:
        return {}
    geometry_by_label: dict[str, dict[str, float]] = {}
    for placeholder in layout_entry.get("placeholders", []):
        label = placeholder.get("label")
        if not label:
            continue
        geometry = dict(placeholder.get("geometry") or {})
        if "fontSizePt" in placeholder:
            geometry["fontSizePt"] = float(placeholder["fontSizePt"])
        if geometry:
            geometry_by_label[str(label)] = geometry
    return geometry_by_label


def _resolve_image_paths(
    sandbox: AgentFileSandbox,
    slide_spec: dict[str, Any],
) -> list[Path]:
    """Resolve sandbox-relative image paths to absolute paths."""

    paths: list[Path] = []
    for image in slide_spec.get("images", []):
        if not isinstance(image, dict):
            continue
        rel = str(image.get("path", "")).strip()
        if not rel:
            continue
        try:
            resolved = sandbox.resolve(rel)
        except AgentFileSandbox.Error as exc:
            raise PptGenerateError("image_load_failed", exc.message) from exc
        if not resolved.is_file():
            raise PptGenerateError("image_load_failed", f"image not found: {rel}")
        if resolved.suffix.lower() not in _IMAGE_SUFFIXES:
            raise PptGenerateError("image_load_failed", f"unsupported image type: {rel}")
        paths.append(resolved)
    return paths


async def _generate_layout_fill(
    outline: dict[str, Any],
    *,
    workspace_id: uuid.UUID,
    output_path: str,
    layout_mode: str,
    chat_model: BaseChatModel | None,
    template_path: Path,
    layout_index_path: Path,
    include_notes: bool,
    transition: str,
) -> dict[str, Any]:
    """Run the layout_fill engine and return the generation result."""

    template = template_path
    layout_json = layout_index_path
    if not template.is_file():
        raise PptGenerateError("template_missing", f"template not found: {template}")

    slide_specs = expand_outline_with_meta(outline)
    layout_index = _load_layout_index(layout_json)
    sandbox = AgentFileSandbox(workspace_id=workspace_id)
    warnings: list[str] = []
    expected_count = len(outline.get("slides", [])) + (1 if outline.get("meta") else 0)
    if len(slide_specs) < expected_count:
        warnings.append("one or more empty outline slides were skipped during normalization")

    prs = Presentation(str(template))
    labels_by_layout = enrich_labels_from_template(
        prs,
        build_placeholder_labels(layout_index),
    )
    remove_all_existing_slides(prs)

    layout_name_to_index = {layout.name: i for i, layout in enumerate(prs.slide_layouts)}
    layout_by_name = {entry["name"]: entry for entry in layout_index}
    layout_by_index = {entry["layoutIndex"]: entry for entry in layout_index}
    pages: list[dict[str, Any]] = []

    for slide_spec in slide_specs:
        layout_name, selection_method, selection_reason = await select_layout_name(
            slide_spec,
            layout_index,
            layout_mode=layout_mode,
            chat_model=chat_model,
        )
        if layout_name not in layout_name_to_index:
            raise PptGenerateError(
                "layout_not_found",
                f"layout not in template: {layout_name}",
            )
        layout_index_no = layout_name_to_index[layout_name]
        slide = prs.slides.add_slide(prs.slide_layouts[layout_index_no])
        labels = labels_by_layout.get(layout_index_no, {})
        layout_entry = layout_by_name.get(layout_name) or layout_by_index.get(layout_index_no)
        layout_ph_geometry = _build_placeholder_geometry(layout_entry)

        image_paths = _resolve_image_paths(sandbox, slide_spec)
        page_warnings = fill_slide_content(
            slide,
            slide_spec,
            labels,
            image_paths=image_paths,
            layout_ph_geometry=layout_ph_geometry,
        )
        warnings.extend(page_warnings)

        if include_notes:
            set_speaker_notes(slide, str(slide_spec.get("speakerNotes", "")))

        if transition != "keep":
            apply_slide_transition(slide, transition)

        pages.append(
            {
                "pageTitle": slide_spec.get("pageTitle"),
                "selectedLayout": layout_name,
                "layoutIndex": layout_index_no,
                "selectionMethod": selection_method,
                "reason": selection_reason,
            }
        )

    out_rel = output_path.strip() or "output/presentation.pptx"
    try:
        dest = sandbox.resolve(out_rel)
    except AgentFileSandbox.Error as exc:
        raise PptGenerateError("write_failed", exc.message) from exc
    dest.parent.mkdir(parents=True, exist_ok=True)
    try:
        prs.save(str(dest))
    except OSError as exc:
        raise PptGenerateError("write_failed", str(exc)) from exc

    return {
        "ok": True,
        "engine": "layout_fill",
        "output_path": out_rel.replace("\\", "/"),
        "pages": pages,
        "warnings": warnings,
    }


def _generate_template_fill(
    outline: dict[str, Any],
    *,
    workspace_id: uuid.UUID,
    output_path: str,
    template_path: Path,
    fill_plan_path: Path | None,
    include_notes: bool,
    transition: str,
) -> dict[str, Any]:
    """Run the template_fill engine and return the generation result."""

    if not template_path.is_file():
        raise PptGenerateError("template_missing", f"template not found: {template_path}")

    sandbox = AgentFileSandbox(workspace_id=workspace_id)
    warnings: list[str] = []

    try:
        library = analyze_template(template_path)
    except Exception as exc:
        raise PptGenerateError("analyze_failed", str(exc)) from exc

    if fill_plan_path is not None:
        if not fill_plan_path.is_file():
            raise PptGenerateError("fill_plan_invalid", f"fill plan not found: {fill_plan_path}")
        try:
            plan = json.loads(fill_plan_path.read_text(encoding="utf-8"))
        except (OSError, json.JSONDecodeError) as exc:
            raise PptGenerateError("fill_plan_invalid", str(exc)) from exc
        if not isinstance(plan, dict) or not plan.get("slides"):
            raise PptGenerateError("fill_plan_invalid", "fill plan must contain slides[]")
    else:
        plan = outline_to_fill_plan(
            outline,
            library,
            source_pptx=str(template_path).replace("\\", "/"),
        )

    check_result = check_fill_plan(library, plan)
    warnings.extend(check_result.get("warnings", []))

    out_rel = output_path.strip() or "output/presentation.pptx"
    try:
        dest = sandbox.resolve(out_rel)
    except AgentFileSandbox.Error as exc:
        raise PptGenerateError("write_failed", exc.message) from exc

    try:
        apply_fill_plan(
            template_path,
            plan,
            dest,
            transition=transition if transition else "fade",
        )
    except (OSError, ValueError) as exc:
        raise PptGenerateError("write_failed", str(exc)) from exc

    slide_specs = expand_outline_with_meta(outline)
    pages: list[dict[str, Any]] = []
    for idx, entry in enumerate(plan.get("slides", [])):
        if not isinstance(entry, dict):
            continue
        title = ""
        if idx < len(slide_specs):
            title = str(slide_specs[idx].get("pageTitle", ""))
        pages.append(
            {
                "pageTitle": title,
                "source_slide": entry.get("source_slide"),
                "replacement_count": len(entry.get("replacements", []) or []),
            }
        )

    if not include_notes:
        pass

    return {
        "ok": True,
        "engine": "template_fill",
        "output_path": out_rel.replace("\\", "/"),
        "pages": pages,
        "warnings": warnings,
    }


def _generate_svg_design(
    *,
    workspace_id: uuid.UUID,
    output_path: str,
    project_dir: str,
    svg_dir: Path | None,
    design_spec_path: Path | None,
    spec_lock_path: Path | None,
    transition: str,
) -> dict[str, Any]:
    """Run the svg_design engine and return the generation result."""

    project = (project_dir or "ppt/default").strip().strip("/") or "ppt/default"
    sandbox = AgentFileSandbox(workspace_id=workspace_id)

    def _resolve_optional(rel_path: str) -> Path:
        try:
            return sandbox.resolve(rel_path)
        except AgentFileSandbox.Error as exc:
            raise PptGenerateError("write_failed", exc.message) from exc

    resolved_svg_dir = svg_dir
    if resolved_svg_dir is None:
        resolved_svg_dir = _resolve_optional(f"{project}/svg")

    resolved_design_spec = design_spec_path
    if resolved_design_spec is None:
        resolved_design_spec = _resolve_optional(f"{project}/svg/design_spec.md")

    resolved_spec_lock = spec_lock_path
    if resolved_spec_lock is None:
        resolved_spec_lock = _resolve_optional(f"{project}/svg/spec_lock.md")

    from app.agent.skills.ppt.svg_pipeline.generate import generate_svg_presentation

    return generate_svg_presentation(
        workspace_id=workspace_id,
        output_path=output_path,
        svg_dir=resolved_svg_dir,
        design_spec_path=resolved_design_spec,
        spec_lock_path=resolved_spec_lock,
        transition=transition,
    )


async def generate_presentation(
    outline: dict[str, Any],
    *,
    workspace_id: uuid.UUID,
    output_path: str,
    layout_mode: str = "hybrid",
    chat_model: BaseChatModel | None = None,
    template_path: Path | None = None,
    layout_index_path: Path | None = None,
    engine: str = "layout_fill",
    include_notes: bool = True,
    transition: str = "fade",
    fill_plan_path: Path | None = None,
    project_dir: str = "ppt/default",
    svg_dir: Path | None = None,
    design_spec_path: Path | None = None,
    spec_lock_path: Path | None = None,
) -> dict[str, Any]:
    """Generate a pptx file in the workspace sandbox from an outline dict."""

    engine_mode = (engine or "layout_fill").strip() or "layout_fill"
    template = template_path or DEFAULT_TEMPLATE_PATH

    if engine_mode == "template_fill":
        return _generate_template_fill(
            outline,
            workspace_id=workspace_id,
            output_path=output_path,
            template_path=template,
            fill_plan_path=fill_plan_path,
            include_notes=include_notes,
            transition=transition,
        )

    if engine_mode == "svg_design":
        return _generate_svg_design(
            workspace_id=workspace_id,
            output_path=output_path,
            project_dir=project_dir,
            svg_dir=svg_dir,
            design_spec_path=design_spec_path,
            spec_lock_path=spec_lock_path,
            transition=transition,
        )

    return await _generate_layout_fill(
        outline,
        workspace_id=workspace_id,
        output_path=output_path,
        layout_mode=layout_mode,
        chat_model=chat_model,
        template_path=template,
        layout_index_path=layout_index_path or DEFAULT_LAYOUT_INDEX_PATH,
        include_notes=include_notes,
        transition=transition,
    )
