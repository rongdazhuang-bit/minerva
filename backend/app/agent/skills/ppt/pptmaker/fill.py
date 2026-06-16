"""Placeholder text values and picture insertion."""

from __future__ import annotations

import re
from pathlib import Path
from typing import Any

from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.presentation import Presentation
from pptx.slide import Slide

from app.agent.skills.ppt.pptmaker.constants import TEMPLATE_PLACEHOLDER_LITERALS
from app.agent.skills.ppt.shared.capacity import check_text_overflow

_FILLABLE_LABEL_PREFIXES = (
    "item",
    "col",
    "grid",
    "left_",
    "right_",
    "metric",
    "toc_",
)


def _text_item_field(items: Any, index: int, field: str) -> str:
    """Read title/body from an items entry that may be dict or plain string."""

    if not isinstance(items, list) or index >= len(items):
        return ""
    item = items[index]
    if isinstance(item, str):
        text = item.strip()
        if field == "body":
            return text
        if field == "title":
            return text
        return ""
    if isinstance(item, dict):
        return str(item.get(field, "")).strip()
    return ""


def _metric_field(metrics: Any, index: int, field: str) -> str:
    """Read number/label/desc from a keyNumbers entry."""

    if not isinstance(metrics, list) or index >= len(metrics):
        return ""
    metric = metrics[index]
    if isinstance(metric, dict):
        return str(metric.get(field, "")).strip()
    if isinstance(metric, str):
        return metric.strip() if field == "number" else ""
    return ""


def _image_caption(images: Any, index: int) -> str:
    """Read caption from an images entry that may be dict or path string."""

    if not isinstance(images, list) or index >= len(images):
        return ""
    image = images[index]
    if isinstance(image, dict):
        return str(image.get("caption", "")).strip()
    return ""


def _is_fillable_label(label: str) -> bool:
    """Return whether a placeholder label is mapped by value_for_label."""

    if label in {"title", "toc_title", "subtitle", "body"}:
        return True
    if label.startswith(_FILLABLE_LABEL_PREFIXES):
        return True
    if "caption" in label.lower() or label.endswith("_desc"):
        return True
    return False


def _assign_placeholder_text(
    shape: Any,
    value: str,
    label: str,
    layout_ph_geometry: dict[str, dict[str, float]] | None,
) -> list[str]:
    """Write text into a placeholder and return overflow warnings."""

    if not hasattr(shape, "text_frame"):
        return []
    shape.text = value
    if not value:
        return []
    return _capacity_warnings_for_text(value, label, layout_ph_geometry)


def value_for_label(slide_spec: dict[str, Any], label: str) -> str:
    """Map a placeholder label to text content from slide_spec."""

    if label in {"title", "toc_title"}:
        return slide_spec.get("pageTitle", "")

    if label == "subtitle":
        return slide_spec.get("subtitle", "")

    if label == "body":
        if slide_spec.get("body"):
            return slide_spec.get("body", "")
        items = slide_spec.get("items", [])
        if len(items) == 1:
            return _text_item_field(items, 0, "body")
        return ""

    if label.startswith("left_") or label.startswith("right_"):
        item_no = 0 if label.startswith("left_") else 1
        field = "title" if label.endswith("_title") else "body"
        return _text_item_field(slide_spec.get("items", []), item_no, field)

    if label.startswith("col") and ("_title" in label or "_body" in label):
        col_no = int(label[3]) - 1
        field = "title" if label.endswith("_title") else "body"
        return _text_item_field(slide_spec.get("items", []), col_no, field)

    if label.startswith("grid") and ("_title" in label or "_body" in label):
        grid_no = int(label[4]) - 1
        field = "title" if label.endswith("_title") else "body"
        return _text_item_field(slide_spec.get("items", []), grid_no, field)

    if label.startswith("item") and ("_title" in label or "_body" in label):
        item_no = int(label[4]) - 1
        field = "title" if label.endswith("_title") else "body"
        return _text_item_field(slide_spec.get("items", []), item_no, field)

    if label.startswith("metric"):
        metric_no = int(label[6]) - 1
        metrics = slide_spec.get("keyNumbers", [])
        if metric_no >= len(metrics):
            return ""
        if label.endswith("_number"):
            return _metric_field(metrics, metric_no, "number")
        if label.endswith("_label"):
            return _metric_field(metrics, metric_no, "label")
        if label.endswith("_desc"):
            return _metric_field(metrics, metric_no, "desc")

    if "caption" in label.lower() or label.endswith("_desc"):
        images = slide_spec.get("images", [])
        match = re.search(r"(\d+)", label)
        if match:
            idx = int(match.group(1)) - 1
            return _image_caption(images, idx)
    return ""


def build_placeholder_labels(layout_index: list[dict[str, Any]]) -> dict[int, dict[int, str]]:
    """Build layoutIndex -> placeholder idx -> label from layout_index.json."""

    labels_by_layout: dict[int, dict[int, str]] = {}
    for layout in layout_index:
        labels_by_layout[layout["layoutIndex"]] = {
            ph["idx"]: ph["label"] for ph in layout.get("placeholders", [])
        }
    return labels_by_layout


def enrich_labels_from_template(
    prs: Presentation,
    labels_by_layout: dict[int, dict[int, str]],
) -> dict[int, dict[int, str]]:
    """Merge runtime placeholder names from the pptx template."""

    for layout_index, layout in enumerate(prs.slide_layouts):
        labels_by_layout.setdefault(layout_index, {})
        for placeholder in layout.placeholders:
            labels_by_layout[layout_index].setdefault(
                placeholder.placeholder_format.idx,
                placeholder.name,
            )
    return labels_by_layout


def _capacity_lines(geometry: dict[str, float], font_size_pt: float) -> int:
    """Estimate line count from placeholder height and font size."""

    height_pt = geometry.get("heightPt")
    if not height_pt or font_size_pt <= 0:
        return 1
    return max(1, int(float(height_pt) / (font_size_pt * 1.2)))


def _capacity_warnings_for_text(
    value: str,
    label: str,
    layout_ph_geometry: dict[str, dict[str, float]] | None,
) -> list[str]:
    """Return overflow warnings for filled text when geometry hints exist."""

    if not value or not layout_ph_geometry:
        return []
    geometry = layout_ph_geometry.get(label)
    if not geometry:
        return []
    width_pt = geometry.get("widthPt")
    if not width_pt:
        return []
    font_size_pt = float(geometry.get("fontSizePt", 18.0))
    return check_text_overflow(
        text=value,
        width_pt=float(width_pt),
        font_size_pt=font_size_pt,
        label=label,
        lines=_capacity_lines(geometry, font_size_pt),
    )


def _caption_placeholders(slide: Slide) -> list[tuple[Any, str]]:
    """Collect non-picture placeholders whose labels suggest captions."""

    result: list[tuple[Any, str]] = []
    for shape in slide.placeholders:
        ph_format = shape.placeholder_format
        if ph_format.type == PP_PLACEHOLDER.PICTURE:
            continue
        label = shape.name
        if "caption" in label.lower() or re.search(r"img\d+_desc|image\d+_desc", label, re.I):
            result.append((shape, label))
    return result


def fill_slide_content(
    slide: Slide,
    slide_spec: dict[str, Any],
    labels: dict[int, str],
    *,
    image_paths: list[Path],
    layout_ph_geometry: dict[str, dict[str, float]] | None = None,
) -> list[str]:
    """Fill text and picture placeholders; return warning messages."""

    warnings: list[str] = []
    picture_slots: list[Any] = []
    for shape in slide.placeholders:
        ph_format = shape.placeholder_format
        if ph_format.type == PP_PLACEHOLDER.PICTURE:
            picture_slots.append(shape)

    for index, image_path in enumerate(image_paths):
        if index >= len(picture_slots):
            warnings.append(f"no picture placeholder for image index {index + 1}")
            break
        placeholder = picture_slots[index]
        try:
            placeholder.insert_picture(str(image_path))
        except Exception as exc:
            raise OSError(f"failed to insert image {image_path}: {exc}") from exc

    for shape in slide.placeholders:
        ph_format = shape.placeholder_format
        if ph_format.type == PP_PLACEHOLDER.PICTURE:
            continue
        label = labels.get(ph_format.idx, shape.name)
        if not _is_fillable_label(label):
            continue
        value = value_for_label(slide_spec, label)
        warnings.extend(_assign_placeholder_text(shape, value, label, layout_ph_geometry))

    for shape, label in _caption_placeholders(slide):
        value = value_for_label(slide_spec, label)
        warnings.extend(_assign_placeholder_text(shape, value, label, layout_ph_geometry))

    for shape in slide.placeholders:
        ph_format = shape.placeholder_format
        if ph_format.type == PP_PLACEHOLDER.PICTURE:
            continue
        if not hasattr(shape, "text_frame"):
            continue
        label = labels.get(ph_format.idx, shape.name)
        current = (shape.text_frame.text or "").strip()
        if current in TEMPLATE_PLACEHOLDER_LITERALS:
            shape.text = ""
            warnings.append(f"cleared unfilled template placeholder '{label}'")

    return warnings


def remove_all_existing_slides(prs: Presentation) -> None:
    """Remove all slides from a presentation, keeping slide layouts."""

    slide_id_list = prs.slides._sldIdLst
    for slide_id in list(slide_id_list):
        rel_id = slide_id.rId
        prs.part.drop_rel(rel_id)
        slide_id_list.remove(slide_id)
