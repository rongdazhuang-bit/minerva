from __future__ import annotations

import json
import re
import sys
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER

from app.agent.skills.ppt.shared.capacity import estimate_text_capacity

EMU_PER_INCH = 914400
_ASSETS_DIR = Path(__file__).resolve().parents[1] / "assets"
DEFAULT_TEMPLATE = _ASSETS_DIR / "template.pptx"
DEFAULT_OUTPUT = _ASSETS_DIR / "layout_index.json"


LAYOUT_TYPE_BY_NAME = {
    "标题幻灯片": "cover",
    "自定义版式": "two_columns",
    "2_自定义版式": "text_with_image",
    "1_自定义版式": "four_grid",
    "要点列表": "bullet_list",
    "三列并列": "three_columns",
    "二列对比": "two_columns_compare",
    "四宫格": "four_grid",
    "大数字展示": "metrics",
    "单段叙述": "single_paragraph",
    "引言金句": "quote",
    "左图右文": "image_left_text_right",
    "右图左文": "text_left_image_right",
    "上图下文": "image_top_text_bottom",
    "下图上文": "text_top_image_bottom",
    "小图大量文字": "small_image_large_text",
    "全图文字叠加": "full_image_text_overlay",
    "两图并排": "two_images",
    "三图横排": "three_images",
    "四图宫格": "four_images",
    "大图小图组合": "hero_image_supporting",
    "五图版式": "five_images",
    "六图宫格": "six_images_grid",
    "目录卡片宫格": "toc_grid",
    "章节左侧色块": "section_divider",
    "章节进度指示": "section_progress",
    "六列并列": "six_columns",
}

# 变体版式名 → 细粒度 layoutType（供规则选版与 LLM 摘要）
LAYOUT_VARIANT_TYPES: dict[str, str] = {
    "要点列表-纵向六条": "bullet_list_vertical",
    "要点列表-横向四项": "bullet_list_horizontal",
    "要点列表-双列六条": "bullet_list_two_columns",
    "三列并列-等宽横排": "three_columns_equal",
    "三列并列-上标题下正文": "three_columns_stacked",
    "三列并列-错落分布": "three_columns_staggered",
    "二列对比-左右等分": "two_columns_compare_equal",
    "二列对比-上下排列": "two_columns_compare_vertical",
    "二列对比-左窄右宽": "two_columns_compare_asymmetric",
    "1_大数字展示-横向四项": "metrics_horizontal",
    "2_大数字展示-横向四项": "metrics_horizontal",
    "大数字展示-横向四项": "metrics_horizontal",
    "大数字展示-田字四项": "metrics_grid",
    "大数字展示-三项居中": "metrics_three_center",
    "四宫格-田字分布": "four_grid_tile",
    "四宫格-横向四项": "four_grid_horizontal",
    "四宫格-纵向四项": "four_grid_vertical",
    "左图右文-左右均分": "image_left_text_right_equal",
    "左图右文-小图大文": "image_left_text_right_small_image",
    "右图左文-左右均分": "text_left_image_right_equal",
    "右图左文-小图大文": "text_left_image_right_small_image",
    "上图下文-半图": "image_top_text_bottom_half",
    "上图下文-浅图": "image_top_text_bottom_light",
    "下图上文-半图": "text_top_image_bottom_half",
    "小图大量文字-右上图": "small_image_large_text_corner",
    "全图文字叠加-居中": "full_image_text_overlay_center",
    "两图并排-带说明": "two_images_with_captions",
    "两图并排-无说明": "two_images",
    "三图横排-带说明": "three_images_with_captions",
    "三图横排-无说明": "three_images",
    "四图宫格-带说明": "four_images_grid_with_captions",
    "四图宫格-无说明": "four_images",
    "大图小图组合-左大右二": "hero_image_supporting_2",
    "大图小图组合-左大右三": "hero_image_supporting_3",
    "五图版式-主次左大": "five_images_hero",
    "五图版式-均等上三下二": "five_images",
    "六图宫格-3x2无说明": "six_images_grid",
    "六图宫格-3x2带说明": "six_images_grid_with_captions",
    "单段叙述-通栏正文": "single_paragraph_full",
    "单段叙述-窄栏正文": "single_paragraph_narrow",
    "引言金句-居中": "quote_center",
    "引言金句-左对齐": "quote_left",
}

LAYOUT_DESCRIPTIONS = {
    "标题幻灯片": {
        "description": "封面页版式，适合整份 PPT 的主标题和副标题，不适合承载正文要点。",
        "bestFor": ["封面", "主标题", "副标题", "开场页"],
        "avoidFor": ["正文要点", "数据指标", "多图展示"],
        "selectionRules": "当页面是整份文档封面、只有主标题和副标题时选择。",
    },
    "自定义版式": {
        "description": "左右两栏正文版式，每栏包含一个标题和一段正文，适合两个对象、两个方案或两类信息并排展示。",
        "bestFor": ["二列对比", "两个方案", "两个对象", "左右并列说明"],
        "avoidFor": ["三个及以上并列要点", "大段单正文", "图片内容"],
        "selectionRules": "当页面有 2 个并列要点，且每个要点都有标题和正文时选择。",
    },
    "2_自定义版式": {
        "description": "左侧两组文字、右侧单图的图文版式，适合一张图片配两段说明或两个局部说明。",
        "bestFor": ["单图说明", "产品截图解释", "图片加两组文字", "图文混排"],
        "avoidFor": ["无图片页面", "多图页面", "四个以上要点"],
        "selectionRules": "当页面有 1 张图片，同时有 1-2 个文字要点时选择。",
    },
    "1_自定义版式": {
        "description": "2x2 四宫格正文版式，每格包含标题和正文，适合四个并列维度或四个模块。",
        "bestFor": ["四个维度", "四类能力", "四项原则", "四模块说明"],
        "avoidFor": ["三项并列", "六项并列", "大数字指标"],
        "selectionRules": "当页面有 4 个普通并列要点，且没有图片和指标数字时选择。",
    },
    "要点列表": {
        "description": "纵向要点列表版式，最多支持 6 个标题加正文的条目，适合作为普通正文页或低置信度兜底。",
        "bestFor": ["普通要点列表", "2-6 个条目", "无法明确判断结构的正文页", "兜底版式"],
        "avoidFor": ["强视觉对比页", "大数字指标页", "明确三列或四宫格结构"],
        "selectionRules": "当内容是普通列表、结构意图不明确，或 AI 置信度较低时选择。",
    },
    "三列并列": {
        "description": "三列等权重并列版式，每列包含标题和正文，适合三个并列优势、方向、能力或模块。",
        "bestFor": ["三大优势", "三个方向", "三类能力", "三项并列说明"],
        "avoidFor": ["两个对象对比", "四个模块", "六个短要点", "数字指标"],
        "selectionRules": "当页面恰好有 3 个普通并列要点，且没有图片和指标数字时优先选择。",
    },
    "四宫格": {
        "description": "2x2 宫格版式，每格包含标题和正文，适合四个等权重内容块。",
        "bestFor": ["四个关键方向", "四个模块", "四项原则", "四类风险"],
        "avoidFor": ["三项并列", "六项并列", "单段正文", "图片展示"],
        "selectionRules": "当页面恰好有 4 个普通并列要点，且没有图片和指标数字时优先选择。",
    },
    "大数字展示": {
        "description": "大数字指标展示版式，最多支持 4 个指标，每个指标包含数字、指标名称和指标说明。",
        "bestFor": ["KPI", "阶段成果", "统计数据", "关键数字", "百分比指标"],
        "avoidFor": ["无数字的普通要点", "长段正文", "图片页"],
        "selectionRules": "当多个要点标题是数字、百分比、MW、万元等指标值，或页面主要表达数据成果时选择。",
    },
    "单段叙述": {
        "description": "大段正文叙述版式，适合一个标题配一整段说明文字。",
        "bestFor": ["背景说明", "概述", "政策说明", "一段连续正文", "长文本说明"],
        "avoidFor": ["多个并列要点", "指标数字页", "图片页", "金句页"],
        "selectionRules": "当页面只有一段正文，或者只有一个 item 且标题为空、正文较长时选择。",
    },
    "引言金句": {
        "description": "强调一句核心观点的金句版式，适合短句、口号、章节引言或结论表达。",
        "bestFor": ["一句话观点", "核心理念", "章节引言", "结论金句"],
        "avoidFor": ["长段正文", "多条要点", "数据指标", "图片页"],
        "selectionRules": "当页面只有一个短标题或一句短而强的观点，正文为空或很短时选择。",
    },
    "左图右文": {
        "description": "左侧图片、右侧正文的单图图文版式，适合以图片为主要视觉证据，右侧解释说明。",
        "bestFor": ["产品图说明", "场景图解释", "案例图片", "单图加正文"],
        "avoidFor": ["无图片页面", "两张以上图片", "纯数据指标"],
        "selectionRules": "当页面有 1 张图片，且文字说明较多时选择。",
    },
    "右图左文": {
        "description": "左侧正文、右侧图片的单图图文版式，适合先说明再展示图片，作为左图右文的节奏变化。",
        "bestFor": ["单图加正文", "产品截图说明", "案例说明", "视觉节奏变化"],
        "avoidFor": ["无图片页面", "多图页面", "纯指标页"],
        "selectionRules": "当页面有 1 张图片且需要文字先行，或前后页面已使用左图右文时选择。",
    },
    "上图下文": {
        "description": "上方大图、下方正文说明的版式，适合先展示场景图或截图，再补充简短说明。",
        "bestFor": ["场景图", "截图展示", "上图下文说明", "视觉优先"],
        "avoidFor": ["文字很多的页面", "无图片页面", "多图页面"],
        "selectionRules": "当页面有 1 张横向图片且文字较短时选择。",
    },
    "两图并排": {
        "description": "两张图片左右并排，并分别带说明，适合前后对比、两个方案、两个案例的视觉展示。",
        "bestFor": ["两图对比", "前后对比", "两个案例", "两个方案截图"],
        "avoidFor": ["一张图片", "三张以上图片", "无图片页面"],
        "selectionRules": "当页面有 2 张图片，且两张图片地位相近时选择。",
    },
    "三图横排": {
        "description": "三张图片横向并列，并分别带说明，适合三个场景、三个案例或三个步骤的图片展示。",
        "bestFor": ["三张图片", "三个案例", "三个场景", "三步图示"],
        "avoidFor": ["图片数量不是 3", "文字很长", "无图片页面"],
        "selectionRules": "当页面有 3 张图片且每张图需要简短说明时选择。",
    },
    "四图宫格": {
        "description": "2x2 四图宫格版式，每张图可带简短说明，适合四个案例、四个场景或产品多角度展示。",
        "bestFor": ["四张图片", "四个案例", "产品多角度", "图片合集"],
        "avoidFor": ["图片数量少于 4", "长文本说明", "无图片页面"],
        "selectionRules": "当页面有 4 张图片，且图片重要性相近时选择。",
    },
    "目录卡片宫格": {
        "description": "目录页卡片宫格版式，支持最多 6 个章节条目，每项包含序号和章节标题。",
        "bestFor": ["目录页", "章节列表", "4-6 个章节", "结构预览"],
        "avoidFor": ["正文页", "数据指标页", "图片展示页"],
        "selectionRules": "当页面类型是目录，且有 4-6 个章节时选择。",
    },
    "章节左侧色块": {
        "description": "章节过渡页版式，左侧显示章节序号，右侧显示章节标题和副标题。",
        "bestFor": ["章节页", "章节切换", "正式汇报过渡页"],
        "avoidFor": ["正文页", "目录页", "多要点页面"],
        "selectionRules": "当页面是章节分隔页，且需要正式、稳重的章节切换感时选择。",
    },
    "章节进度指示": {
        "description": "带进度信息的章节页版式，适合章节较多时提示当前位置。",
        "bestFor": ["章节页", "进度提示", "第几章/共几章", "章节定位"],
        "avoidFor": ["正文页", "目录页", "普通要点页"],
        "selectionRules": "当页面是章节页，且内容中包含当前章节进度或需要位置提示时选择。",
    },
    "六列并列": {
        "description": "六项并列版式，采用 3x2 排布，每项包含标题和正文，适合六个等权重指标、因素或说明项。",
        "bestFor": ["六个要点", "六项指标说明", "六个因素", "六类能力", "3x2 并列内容"],
        "avoidFor": ["少于 5 个要点", "大段正文", "数字指标大卡片", "图片内容"],
        "selectionRules": "当页面有 6 个普通并列要点，且每个要点包含标题和正文、没有图片和大数字展示需求时选择。",
    },
    "二列对比": {
        "description": "二列对比版式，适合两个对象、两个方案或两类信息的差异比较。",
        "bestFor": ["两个对象", "两个方案", "现状 vs 目标", "优劣对比"],
        "avoidFor": ["三个及以上要点", "单段正文", "图片展示"],
        "selectionRules": "当页面有 2 个并列要点，且语义上存在对比或左右并列关系时选择。",
    },
    "下图上文": {
        "description": "上方文字、下方图片的单图图文版式，适合先说明后展示。",
        "bestFor": ["先文字后图片", "单图说明", "文字优先"],
        "avoidFor": ["无图片页面", "多图页面"],
        "selectionRules": "当页面有 1 张图片且需要文字先行时选择。",
    },
    "小图大量文字": {
        "description": "小图配大量文字的图文版式，图片作为辅助说明，正文为主要内容。",
        "bestFor": ["文字较多", "辅助图片", "说明型图文页"],
        "avoidFor": ["图片为主的展示页", "多图页面"],
        "selectionRules": "当页面有 1 张图片但文字明显多于图片说明时选择。",
    },
    "全图文字叠加": {
        "description": "全页图片加文字叠加版式，适合视觉冲击型单图页面。",
        "bestFor": ["全图展示", "视觉冲击", "图片背景加短文案"],
        "avoidFor": ["长正文", "多图页面", "需要强可读性的密集内容"],
        "selectionRules": "当页面有 1 张主视觉图片，且叠加文字很短时选择。",
    },
    "大图小图组合": {
        "description": "一张主图搭配多张小图的组合版式，适合主次关系明确的图片展示。",
        "bestFor": ["主图加细节图", "主场景加局部", "有主次的多图"],
        "avoidFor": ["图片地位完全均等", "无图片页面"],
        "selectionRules": "当页面有 3-4 张图片，且其中一张明显更重要时选择。",
    },
    "五图版式": {
        "description": "五张图片展示版式，支持主次型或均等型五图排布。",
        "bestFor": ["五张图片", "产品系列", "多场景展示"],
        "avoidFor": ["图片少于 5", "长文本说明"],
        "selectionRules": "当页面有 5 张图片时选择，并根据是否有主次选择具体变体。",
    },
    "六图宫格": {
        "description": "六张图片的 3x2 宫格版式，可无说明或每图带简短说明。",
        "bestFor": ["六张图片", "案例集", "图片墙", "产品系列"],
        "avoidFor": ["图片少于 6", "长文本说明"],
        "selectionRules": "当页面有 6 张图片，且图片地位相近时选择。",
    },
}


def placeholder_type_name(ph_type: PP_PLACEHOLDER) -> str:
    if ph_type == PP_PLACEHOLDER.TITLE or ph_type == PP_PLACEHOLDER.CENTER_TITLE:
        return "title"
    if ph_type == PP_PLACEHOLDER.SUBTITLE:
        return "subtitle"
    if ph_type == PP_PLACEHOLDER.BODY:
        return "body"
    if ph_type == PP_PLACEHOLDER.PICTURE:
        return "picture"
    return str(ph_type).split(" ")[0].lower()


def inches(value: int) -> float:
    """Convert EMU to inches rounded to three decimals."""

    return round(value / EMU_PER_INCH, 3)


def points_from_emu(value: int) -> float:
    """Convert EMU to typographic points."""

    return round(value / EMU_PER_INCH * 72, 1)


def placeholder_geometry(shape) -> dict[str, float]:
    """Return width/height in points for capacity estimation."""

    return {
        "widthPt": points_from_emu(shape.width),
        "heightPt": points_from_emu(shape.height),
    }


def placeholder_font_size_pt(shape) -> float | None:
    """Read the first run font size in points when a text frame exists."""

    if not shape.has_text_frame:
        return None
    text_frame = shape.text_frame
    if not text_frame.paragraphs:
        return None
    paragraph = text_frame.paragraphs[0]
    if not paragraph.runs:
        return None
    font = paragraph.runs[0].font
    if font.size is None:
        return None
    return round(font.size.pt, 1)


def label_index_count(labels: list[str], prefix: str, suffixes: tuple[str, ...]) -> int:
    found = set()
    pattern = re.compile(rf"^{re.escape(prefix)}(\d+)_({'|'.join(map(re.escape, suffixes))})$")
    for label in labels:
        match = pattern.match(label)
        if match:
            found.add(int(match.group(1)))
    return len(found)


def resolve_layout_type(name: str) -> str:
    """Resolve base or variant layout name to a layoutType string."""

    if name in LAYOUT_VARIANT_TYPES:
        return LAYOUT_VARIANT_TYPES[name]
    if name in LAYOUT_TYPE_BY_NAME:
        return LAYOUT_TYPE_BY_NAME[name]
    base_name = re.sub(r"^\d+_", "", name.split("-", 1)[0])
    return LAYOUT_TYPE_BY_NAME.get(base_name, "unknown")


def compute_capacity_hints(placeholders: list[dict[str, Any]]) -> dict[str, Any]:
    """Aggregate per-label text capacity estimates for layout-first selection."""

    by_label: dict[str, dict[str, float | int]] = {}
    body_capacities: list[float] = []
    title_capacities: list[float] = []

    for placeholder in placeholders:
        label = str(placeholder.get("label", ""))
        geometry = placeholder.get("geometry") or {}
        width_pt = geometry.get("widthPt")
        if not width_pt:
            continue
        height_pt = float(geometry.get("heightPt") or 0)
        font_size_pt = float(placeholder.get("fontSizePt") or 18.0)
        lines = max(1, int(height_pt / (font_size_pt * 1.2))) if height_pt else 1
        capacity = estimate_text_capacity(
            width_pt=float(width_pt),
            font_size_pt=font_size_pt,
            lines=lines,
        )
        by_label[label] = {
            "widthPt": round(float(width_pt), 1),
            "heightPt": round(height_pt, 1),
            "fontSizePt": font_size_pt,
            "lines": lines,
            "capacityUnits": round(capacity, 1),
        }
        if label == "body" or label.endswith("_body"):
            body_capacities.append(capacity)
        if label == "title" or label.endswith("_title"):
            title_capacities.append(capacity)

    return {
        "byLabel": by_label,
        "minBodyCapacity": round(min(body_capacities), 1) if body_capacities else None,
        "maxBodyCapacity": round(max(body_capacities), 1) if body_capacities else None,
        "avgBodyCapacity": round(sum(body_capacities) / len(body_capacities), 1)
        if body_capacities
        else None,
        "minTitleCapacity": round(min(title_capacities), 1) if title_capacities else None,
    }


def infer_page_type(name: str, layout_type: str) -> str:
    if layout_type == "cover":
        return "cover"
    if layout_type.startswith("toc"):
        return "toc"
    if layout_type.startswith("section"):
        return "chapter"
    return "content"


def infer_content_signals(name: str, placeholders: list[dict[str, Any]]) -> dict[str, Any]:
    labels = [item["label"] for item in placeholders]
    layout_type = resolve_layout_type(name)
    has_image = any(item["type"] == "picture" for item in placeholders)
    image_count = sum(1 for item in placeholders if item["type"] == "picture")

    item_slots = label_index_count(labels, "item", ("title", "body"))
    col_slots = label_index_count(labels, "col", ("title", "body"))
    grid_slots = label_index_count(labels, "grid", ("title", "body"))
    metric_slots = label_index_count(labels, "metric", ("number", "label", "desc"))
    toc_slots = 0
    for label in labels:
        match = re.match(r"^toc_title_(\d+)$", label)
        if match:
            toc_slots = max(toc_slots, int(match.group(1)))

    if "left_title" in labels and "right_title" in labels:
        item_slots = max(item_slots, 2)
    if layout_type in {"single_paragraph", "quote"}:
        item_slots = max(item_slots, 1)

    max_items = max(item_slots, col_slots, grid_slots, metric_slots, toc_slots)
    if max_items == 0 and image_count:
        max_items = image_count

    columns_by_type = {
        "cover": 1,
        "two_columns": 2,
        "two_columns_compare": 2,
        "text_with_image": 2,
        "three_columns": 3,
        "four_grid": 2,
        "six_columns": 6,
        "metrics": metric_slots or 4,
        "single_paragraph": 1,
        "quote": 1,
        "image_left_text_right": 2,
        "text_left_image_right": 2,
        "image_top_text_bottom": 1,
        "text_top_image_bottom": 1,
        "small_image_large_text": 2,
        "full_image_text_overlay": 1,
        "two_images": 2,
        "three_images": 3,
        "four_images": 2,
        "hero_image_supporting": 2,
        "five_images": 3,
        "six_images_grid": 3,
        "toc_grid": 3,
        "section_divider": 2,
        "section_progress": 1,
        "bullet_list": 1,
    }

    columns = columns_by_type.get(layout_type)
    if columns is None:
        for base_type in sorted(columns_by_type, key=len, reverse=True):
            if layout_type == base_type or layout_type.startswith(f"{base_type}_"):
                columns = columns_by_type[base_type]
                break
    if columns is None:
        columns = 1
    if layout_type.startswith("metrics") and metric_slots:
        columns = metric_slots

    return {
        "pageType": infer_page_type(name, layout_type),
        "layoutType": layout_type,
        "columns": columns,
        "hasImage": has_image,
        "imageCount": image_count,
        "maxItems": max_items,
        "supportsTitle": any(item["label"] == "title" for item in placeholders),
        "supportsBody": any(item["label"] == "body" for item in placeholders),
        "supportsMetrics": metric_slots > 0,
        "supportsQuote": "quote" in labels,
    }


def describe_layout(name: str) -> dict[str, Any]:
    if name in LAYOUT_DESCRIPTIONS:
        return LAYOUT_DESCRIPTIONS[name]
    base_name = re.sub(r"^\d+_", "", name.split("-", 1)[0])
    if base_name in LAYOUT_DESCRIPTIONS:
        base = dict(LAYOUT_DESCRIPTIONS[base_name])
        base["description"] = f"{base['description']} 当前变体：{name}。"
        return base
    return {
        "description": "未配置语义描述的版式，请优先根据占位符、列数、图片数和名称判断。",
        "bestFor": [],
        "avoidFor": [],
        "selectionRules": "当没有更明确匹配时才选择。",
    }


def shape_distribution(layout) -> dict[str, Any]:
    placeholders = list(layout.placeholders)
    static_count = len(layout.shapes) - len(placeholders)
    xs = sorted({round(inches(shape.left), 1) for shape in placeholders})
    ys = sorted({round(inches(shape.top), 1) for shape in placeholders})
    return {
        "shapeCount": len(layout.shapes),
        "placeholderCount": len(placeholders),
        "staticShapeCount": static_count,
        "distinctXPositions": xs,
        "distinctYPositions": ys,
    }


def extract_layout_index(template_path: Path) -> list[dict[str, Any]]:
    prs = Presentation(str(template_path))
    layout_index = []

    for layout_no, layout in enumerate(prs.slide_layouts):
        placeholders = []
        for placeholder in layout.placeholders:
            fmt = placeholder.placeholder_format
            entry: dict[str, Any] = {
                "idx": fmt.idx,
                "type": placeholder_type_name(fmt.type),
                "label": placeholder.name,
                "x": inches(placeholder.left),
                "y": inches(placeholder.top),
                "w": inches(placeholder.width),
                "h": inches(placeholder.height),
                "geometry": placeholder_geometry(placeholder),
            }
            font_size_pt = placeholder_font_size_pt(placeholder)
            if font_size_pt is not None:
                entry["fontSizePt"] = font_size_pt
            placeholders.append(entry)

        placeholders.sort(key=lambda item: (item["idx"], item["y"], item["x"]))

        layout_index.append({
            "layoutIndex": layout_no,
            "name": layout.name,
            **describe_layout(layout.name),
            "sourcePart": f"ppt/slideLayouts/slideLayout{layout_no + 1}.xml",
            "placeholders": placeholders,
            "contentSignals": infer_content_signals(layout.name, placeholders),
            "capacityHints": compute_capacity_hints(placeholders),
            "shapeSignals": shape_distribution(layout),
        })

    return layout_index


def main() -> None:
    template_path = Path(sys.argv[1]) if len(sys.argv) > 1 else DEFAULT_TEMPLATE
    output_path = Path(sys.argv[2]) if len(sys.argv) > 2 else DEFAULT_OUTPUT
    index = extract_layout_index(template_path)
    output_path.write_text(
        json.dumps(index, ensure_ascii=False, indent=2),
        encoding="utf-8",
    )
    print(f"layouts: {len(index)}")
    print(f"saved: {output_path}")


if __name__ == "__main__":
    main()
