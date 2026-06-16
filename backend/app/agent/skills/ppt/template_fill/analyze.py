"""Extract slide library schema from a source PPTX template."""

from __future__ import annotations

import re
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Emu

_SCHEMA = "template_fill_pptx_library.v1"
_EMU_PER_PX = 914400 / 96
_COVER_KEYWORDS = ("封面", "cover", "title slide", "主标题")
_SECTION_KEYWORDS = ("章节", "section", "part", "篇章")
_CLOSING_KEYWORDS = ("谢谢", "thank", "q&a", "questions", "结语", "结束")

_ROLE_TITLE = "title_candidate"
_ROLE_BODY = "body_candidate"
_ROLE_LABEL = "label_candidate"


def _emu_to_px(value: int | Emu | None) -> float:
    """Convert EMU to approximate pixel units at 96 DPI."""

    if value is None:
        return 0.0
    return float(int(value)) / _EMU_PER_PX


def _pt_to_px(pt: float) -> float:
    """Convert typographic points to approximate pixels."""

    return pt * (96.0 / 72.0)


def _shape_font_size_px(shape) -> float:
    """Estimate the dominant font size in points for a text shape."""

    if not getattr(shape, "has_text_frame", False):
        return 18.0
    sizes: list[float] = []
    for paragraph in shape.text_frame.paragraphs:
        if paragraph.font.size is not None:
            sizes.append(float(paragraph.font.size.pt))
        for run in paragraph.runs:
            if run.font.size is not None:
                sizes.append(float(run.font.size.pt))
    if sizes:
        return _pt_to_px(max(sizes))
    return _pt_to_px(18.0)


def _shape_text(shape) -> str:
    """Return trimmed text from a shape when it has a text frame."""

    if not getattr(shape, "has_text_frame", False):
        return ""
    return (shape.text_frame.text or "").strip()


def _infer_role(shape, *, area: float, slide_height: float) -> str:
    """Infer slot role from geometry, font size, and text length."""

    text = _shape_text(shape)
    if not text:
        return _ROLE_LABEL
    font_px = _shape_font_size_px(shape)
    top = _emu_to_px(shape.top)
    text_len = len(text)
    if top < slide_height * 0.25 and font_px >= 24 and text_len <= 80:
        return _ROLE_TITLE
    if font_px <= 14 or text_len <= 12:
        return _ROLE_LABEL
    return _ROLE_BODY


def _infer_page_type(
    slide_index: int,
    slide_count: int,
    slots: list[dict[str, Any]],
) -> str:
    """Infer page_type from slide index and aggregated slot text."""

    combined = " ".join(str(slot.get("text", "")) for slot in slots).lower()
    if slide_index == 1:
        return "cover_candidate"
    if slide_index == slide_count and any(kw in combined for kw in _CLOSING_KEYWORDS):
        return "closing_candidate"
    if any(kw in combined for kw in _SECTION_KEYWORDS):
        return "section_candidate"
    if any(kw in combined for kw in _COVER_KEYWORDS) and len(slots) <= 3:
        return "cover_candidate"
    title_slots = [s for s in slots if s.get("role") == _ROLE_TITLE]
    body_slots = [s for s in slots if s.get("role") == _ROLE_BODY]
    if title_slots and not body_slots and len(slots) <= 2:
        return "section_candidate"
    return "content_candidate"


def _iter_text_shapes(slide):
    """Yield shapes that expose a text frame, including grouped children."""

    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            for child in shape.shapes:
                if getattr(child, "has_text_frame", False):
                    yield child
        elif getattr(shape, "has_text_frame", False):
            yield shape


def _build_slot(slide_index: int, shape, slide_height: float) -> dict[str, Any]:
    """Build one slot entry for a text shape."""

    text = _shape_text(shape)
    width = _emu_to_px(shape.width)
    height = _emu_to_px(shape.height)
    area = width * height
    role = _infer_role(shape, area=area, slide_height=slide_height)
    paragraph_count = 0
    if getattr(shape, "has_text_frame", False):
        paragraph_count = len([p for p in shape.text_frame.paragraphs if (p.text or "").strip()])
    return {
        "slot_id": f"s{slide_index:02d}_sh{shape.shape_id}",
        "role": role,
        "text": text,
        "geometry": {
            "x": round(_emu_to_px(shape.left), 1),
            "y": round(_emu_to_px(shape.top), 1),
            "w": round(width, 1),
            "h": round(height, 1),
        },
        "text_metrics": {
            "font_size_px": round(_shape_font_size_px(shape), 1),
            "paragraph_count": max(paragraph_count, 1 if text else 0),
        },
    }


def analyze_template(pptx_path: Path) -> dict[str, Any]:
    """Return slide library schema ``template_fill_pptx_library.v1`` for a source deck."""

    prs = Presentation(str(pptx_path))
    slide_height = _emu_to_px(prs.slide_height)
    slides_out: list[dict[str, Any]] = []

    for idx, slide in enumerate(prs.slides, start=1):
        slots = [_build_slot(idx, shape, slide_height) for shape in _iter_text_shapes(slide)]
        page_type = _infer_page_type(idx, len(prs.slides), slots)
        text_summary = " | ".join(s["text"][:60] for s in slots if s["text"])[:240]
        slides_out.append(
            {
                "slide_index": idx,
                "page_type": page_type,
                "text_summary": text_summary,
                "slots": slots,
                "tables": [],
                "charts": [],
            }
        )

    return {
        "schema": _SCHEMA,
        "source_pptx": str(pptx_path).replace("\\", "/"),
        "slide_count": len(slides_out),
        "slides": slides_out,
    }


def parse_shape_id_from_slot_id(slot_id: str) -> int | None:
    """Extract shape_id integer suffix from a slot_id like ``s01_sh42``."""

    match = re.search(r"_sh(\d+)$", slot_id)
    if not match:
        return None
    return int(match.group(1))
