"""Apply a fill plan to a source PPTX and write the output deck."""

from __future__ import annotations

import copy
from pathlib import Path
from typing import Any

from pptx import Presentation

from app.agent.skills.ppt.shared.notes import set_speaker_notes
from app.agent.skills.ppt.shared.transitions import apply_slide_transition
from app.agent.skills.ppt.template_fill.analyze import parse_shape_id_from_slot_id


def _blank_layout(prs: Presentation):
    """Return the most blank-like slide layout available."""

    for layout in prs.slide_layouts:
        if layout.name.lower() in {"blank", "空白", "empty"}:
            return layout
    return prs.slide_layouts[-1]


def _clone_slide(source_prs: Presentation, source_index: int, dest_prs: Presentation):
    """Deep-copy one slide from source_prs into dest_prs."""

    source_slide = source_prs.slides[source_index]
    dest_slide = dest_prs.slides.add_slide(_blank_layout(dest_prs))

    for shape in source_slide.shapes:
        new_element = copy.deepcopy(shape.element)
        dest_slide.shapes._spTree.insert_element_before(new_element, "p:extLst")

    return dest_slide


def _find_shape_by_id(slide, shape_id: int):
    """Find a shape on slide by python-pptx shape_id."""

    for shape in slide.shapes:
        if shape.shape_id == shape_id:
            return shape
    return None


def _set_shape_text(shape, text: str) -> None:
    """Replace all text in a shape text frame."""

    if not getattr(shape, "has_text_frame", False):
        return
    shape.text_frame.clear()
    shape.text_frame.text = text


def _apply_replacements(slide, replacements: list[dict[str, Any]]) -> None:
    """Apply slot replacements to shapes on a cloned slide."""

    for replacement in replacements:
        if not isinstance(replacement, dict):
            continue
        slot_id = str(replacement.get("slot_id", ""))
        text = str(replacement.get("text", ""))
        shape_id = parse_shape_id_from_slot_id(slot_id)
        if shape_id is None:
            continue
        shape = _find_shape_by_id(slide, shape_id)
        if shape is not None:
            _set_shape_text(shape, text)


def apply_fill_plan(
    source_pptx: Path,
    plan: dict[str, Any],
    output_pptx: Path,
    *,
    transition: str = "fade",
) -> None:
    """Clone selected source slides, apply replacements, notes, and transitions."""

    source_prs = Presentation(str(source_pptx))
    dest_prs = Presentation()
    dest_prs.slide_width = source_prs.slide_width
    dest_prs.slide_height = source_prs.slide_height

    plan_slides = plan.get("slides", [])
    if not isinstance(plan_slides, list):
        raise ValueError("fill plan slides must be a list")

    for entry in plan_slides:
        if not isinstance(entry, dict):
            continue
        source_slide = int(entry.get("source_slide", 0))
        if source_slide < 1 or source_slide > len(source_prs.slides):
            raise ValueError(f"invalid source_slide: {source_slide}")
        cloned = _clone_slide(source_prs, source_slide - 1, dest_prs)
        replacements = entry.get("replacements", [])
        if isinstance(replacements, list):
            _apply_replacements(cloned, replacements)
        notes = str(entry.get("notes", "")).strip()
        if notes:
            set_speaker_notes(cloned, notes)
        if transition != "keep":
            apply_slide_transition(cloned, transition)

    output_pptx.parent.mkdir(parents=True, exist_ok=True)
    dest_prs.save(str(output_pptx))
