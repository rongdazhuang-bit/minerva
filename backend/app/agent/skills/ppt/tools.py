"""PPT skill tools: ingest, analyze, check, validate, draft, and generate."""

from __future__ import annotations

import asyncio
import json
import uuid
from pathlib import Path
from typing import Any

from langchain_core.messages import HumanMessage, SystemMessage
from langchain_core.tools import tool
from pptx import Presentation

from app.agent.infrastructure.agent_file_sandbox import AgentFileSandbox
from app.agent.infrastructure.skill_tool_context import SkillToolContext
from app.agent.skills.ppt.ingest.converters import (
    IngestError,
    build_image_manifest,
    convert_file_to_markdown,
    convert_url_to_markdown,
)
from app.agent.skills.ppt.pptmaker.constants import DEFAULT_TEMPLATE_PATH
from app.agent.skills.ppt.pptmaker.generate import PptGenerateError, generate_presentation
from app.agent.skills.ppt.pptmaker.layout_select import extract_json_object
from app.agent.skills.ppt.pptmaker.schemas import validate_outline_dict
from app.agent.skills.ppt.template_fill.analyze import analyze_template
from app.agent.skills.ppt.template_fill.check_plan import check_fill_plan

_OUTLINE_SYSTEM = """你是 PPT 大纲撰写助手。根据用户 brief 生成 JSON 大纲，且只输出 JSON，不要 markdown 围栏。

输出 schema：
{
  "meta": { "title": "<封面主标题>", "subtitle": "<封面副标题>" },
  "slides": [
    {
      "pageTitle": "<本页标题>",
      "pageType": "toc",
      "items": ["<目录条目1>", "<目录条目2>"] 或 [{ "title": "<要点小标题>", "body": "<要点说明>" }],
      "keyNumbers": [{ "number": "19", "label": "<指标名>", "desc": "<指标说明>" }],
      "body": "<单段长正文>",
      "images": [{ "path": "<沙箱相对路径>", "caption": "<图注>" }],
      "hasImage": false,
      "speakerNotes": "<演讲者备注，可选>"
    }
  ]
}

规则：
- meta 可选；有封面需求时填写 title/subtitle，不要再在 slides 里重复 pageType=cover。
- 每页只选一种主要内容：items 要点、keyNumbers 指标、body 长正文、或 images 配图。
- 目录页：pageTitle 用「目录」，items 写字符串数组（每章一行），或设置 pageType=toc。
- 禁止使用 schema 尖括号外的示例占位字面量（如「要点标题」「要点正文」）；必须写真实内容。
- images[].path 必须是用户提供的沙箱相对路径；无图片时不要写 images。
- 指标页用 keyNumbers，数字放 number 字段。
- 不要输出空 slides 或无内容的页面。"""


async def _sandbox_write_json(
    workspace_id: uuid.UUID,
    path: str,
    data: dict[str, Any],
) -> dict[str, object]:
    """Write JSON to the workspace sandbox."""

    box = AgentFileSandbox(workspace_id=workspace_id)
    content = json.dumps(data, ensure_ascii=False, indent=2)
    return await box.write_file_async(path, content)


def _collect_slide_texts(pptx_path: Path) -> list[str]:
    """Collect non-empty shape texts from all slides in a pptx."""

    prs = Presentation(str(pptx_path))
    texts: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                text = (shape.text_frame.text or "").strip()
                if text:
                    texts.append(text)
    return texts


def validate_ppt_output(
    pptx_path: Path,
    *,
    expected_slide_count: int | None = None,
    expected_texts: list[str] | None = None,
) -> dict[str, Any]:
    """Read a pptx and validate slide count and expected texts."""

    checks: list[dict[str, Any]] = []
    warnings: list[str] = []

    if not pptx_path.is_file():
        return {
            "passed": False,
            "checks": [{"name": "file_exists", "passed": False}],
            "warnings": [f"pptx not found: {pptx_path}"],
        }

    prs = Presentation(str(pptx_path))
    slide_count = len(prs.slides)
    texts = _collect_slide_texts(pptx_path)
    combined = "\n".join(texts)

    if expected_slide_count is not None:
        passed = slide_count == expected_slide_count
        checks.append(
            {
                "name": "slide_count",
                "passed": passed,
                "expected": expected_slide_count,
                "actual": slide_count,
            }
        )
        if not passed:
            warnings.append(f"expected {expected_slide_count} slides, got {slide_count}")

    if expected_texts:
        for expected in expected_texts:
            passed = expected in combined
            checks.append(
                {
                    "name": "expected_text",
                    "passed": passed,
                    "expected": expected,
                }
            )
            if not passed:
                warnings.append(f"expected text not found: {expected}")

    passed = all(item.get("passed", False) for item in checks) if checks else True
    return {"passed": passed, "checks": checks, "warnings": warnings, "slide_count": slide_count}


def register_tools(ctx: SkillToolContext) -> list[Any]:
    """Register PPT maker tools bound to the current workspace and chat model."""

    workspace_id = ctx.workspace_id
    chat_model = ctx.chat_model

    @tool
    async def ingest_ppt_source(
        source_path: str = "",
        url: str = "",
        project_dir: str = "ppt/default",
        output_md_path: str = "",
    ) -> str:
        """将沙箱内源文件或 URL 转为 Markdown，写入项目 sources 目录。"""

        project = (project_dir or "ppt/default").strip().strip("/") or "ppt/default"
        md_rel = (output_md_path or f"{project}/sources/content.md").strip()
        images_rel = f"{project}/sources/images"
        manifest_rel = f"{project}/sources/image_manifest.json"
        warnings: list[str] = []

        src = (source_path or "").strip()
        target_url = (url or "").strip()
        if not src and not target_url:
            return json.dumps(
                {
                    "ok": False,
                    "error": "source_path or url is required",
                    "code": "source_missing",
                },
                ensure_ascii=False,
            )

        box = AgentFileSandbox(workspace_id=workspace_id)
        image_paths: list[Path] = []

        try:
            if src:
                local_source = box.resolve(src)
                if not local_source.is_file():
                    return json.dumps(
                        {
                            "ok": False,
                            "error": f"source not found: {src}",
                            "code": "source_missing",
                        },
                        ensure_ascii=False,
                    )
                images_local = box.resolve(images_rel)
                markdown, image_paths = await asyncio.to_thread(
                    convert_file_to_markdown,
                    local_source,
                    images_dir=images_local,
                )
            else:
                markdown, image_paths = await asyncio.to_thread(
                    convert_url_to_markdown,
                    target_url,
                )
        except IngestError as exc:
            return json.dumps(
                {"ok": False, "error": exc.message, "code": exc.code},
                ensure_ascii=False,
            )
        except Exception as exc:
            return json.dumps(
                {"ok": False, "error": str(exc), "code": "convert_failed"},
                ensure_ascii=False,
            )

        try:
            write_result = await box.write_file_async(md_rel, markdown)
            if not write_result.get("ok"):
                return json.dumps(write_result, ensure_ascii=False)

            manifest_base = box.workspace_root()
            manifest = build_image_manifest(image_paths, manifest_base)
            manifest_write = await box.write_file_async(
                manifest_rel,
                json.dumps(manifest, ensure_ascii=False, indent=2),
            )
            if not manifest_write.get("ok"):
                return json.dumps(manifest_write, ensure_ascii=False)
        except AgentFileSandbox.Error as exc:
            return json.dumps(
                {"ok": False, "error": exc.message, "code": "write_failed"},
                ensure_ascii=False,
            )

        return json.dumps(
            {
                "ok": True,
                "md_path": str(write_result.get("path", md_rel)),
                "images_dir": images_rel,
                "manifest_path": str(manifest_write.get("path", manifest_rel)),
                "warnings": warnings,
            },
            ensure_ascii=False,
        )

    @tool
    async def analyze_ppt_template(
        template_path: str,
        output_path: str = "",
        project_dir: str = "ppt/default",
    ) -> str:
        """分析沙箱内 PPTX 模板，输出 slide_library.json。"""

        project = (project_dir or "ppt/default").strip().strip("/") or "ppt/default"
        out_rel = (output_path or f"{project}/analysis/slide_library.json").strip()
        rel_template = (template_path or "").strip()
        if not rel_template:
            return json.dumps(
                {"ok": False, "error": "template_path is required", "code": "template_missing"},
                ensure_ascii=False,
            )

        box = AgentFileSandbox(workspace_id=workspace_id)
        try:
            local_template = box.resolve(rel_template)
        except AgentFileSandbox.Error as exc:
            return json.dumps(
                {"ok": False, "error": exc.message, "code": "template_missing"},
                ensure_ascii=False,
            )
        if not local_template.is_file():
            return json.dumps(
                {"ok": False, "error": f"template not found: {rel_template}", "code": "template_missing"},
                ensure_ascii=False,
            )

        try:
            library = await asyncio.to_thread(analyze_template, local_template)
        except Exception as exc:
            return json.dumps(
                {"ok": False, "error": str(exc), "code": "analyze_failed"},
                ensure_ascii=False,
            )

        try:
            write_result = await _sandbox_write_json(workspace_id, out_rel, library)
        except AgentFileSandbox.Error as exc:
            return json.dumps(
                {"ok": False, "error": exc.message, "code": "write_failed"},
                ensure_ascii=False,
            )
        if not write_result.get("ok"):
            return json.dumps(write_result, ensure_ascii=False)

        page_types = [s.get("page_type") for s in library.get("slides", []) if isinstance(s, dict)]
        return json.dumps(
            {
                "ok": True,
                "path": str(write_result.get("path", out_rel)),
                "slide_count": library.get("slide_count", 0),
                "page_types": page_types,
            },
            ensure_ascii=False,
        )

    @tool
    async def check_ppt_fill_plan(
        slide_library_path: str,
        fill_plan_path: str,
        output_path: str = "",
        project_dir: str = "ppt/default",
    ) -> str:
        """校验 fill plan 与 slide library 的容量与 slot 一致性。"""

        project = (project_dir or "ppt/default").strip().strip("/") or "ppt/default"
        out_rel = (output_path or f"{project}/analysis/check_report.json").strip()
        box = AgentFileSandbox(workspace_id=workspace_id)

        for label, rel_path, code in (
            ("slide_library_path", slide_library_path, "fill_plan_invalid"),
            ("fill_plan_path", fill_plan_path, "fill_plan_invalid"),
        ):
            rel = (rel_path or "").strip()
            if not rel:
                return json.dumps(
                    {"ok": False, "error": f"{label} is required", "code": code},
                    ensure_ascii=False,
                )
            try:
                read_result = await box.read_file_async(rel)
            except AgentFileSandbox.Error as exc:
                return json.dumps(
                    {"ok": False, "error": exc.message, "code": code},
                    ensure_ascii=False,
                )
            if not read_result.get("ok"):
                return json.dumps(read_result, ensure_ascii=False)
            if label == "slide_library_path":
                try:
                    library = json.loads(str(read_result.get("content", "")))
                except json.JSONDecodeError as exc:
                    return json.dumps(
                        {"ok": False, "error": str(exc), "code": "fill_plan_invalid"},
                        ensure_ascii=False,
                    )
            else:
                try:
                    plan = json.loads(str(read_result.get("content", "")))
                except json.JSONDecodeError as exc:
                    return json.dumps(
                        {"ok": False, "error": str(exc), "code": "fill_plan_invalid"},
                        ensure_ascii=False,
                    )

        report = check_fill_plan(library, plan)
        try:
            write_result = await _sandbox_write_json(workspace_id, out_rel, report)
        except AgentFileSandbox.Error as exc:
            return json.dumps(
                {"ok": False, "error": exc.message, "code": "write_failed"},
                ensure_ascii=False,
            )
        if not write_result.get("ok"):
            return json.dumps(write_result, ensure_ascii=False)

        return json.dumps(
            {
                "ok": True,
                "path": str(write_result.get("path", out_rel)),
                "passed": report.get("passed", False),
                "warnings": report.get("warnings", []),
            },
            ensure_ascii=False,
        )

    @tool
    async def validate_ppt_output_tool(
        pptx_path: str,
        expected_slide_count: int | None = None,
        expected_texts: str = "",
    ) -> str:
        """回读 pptx 校验页数与关键文本是否存在。"""

        rel = (pptx_path or "").strip()
        if not rel:
            return json.dumps(
                {"ok": False, "error": "pptx_path is required", "code": "validation_error"},
                ensure_ascii=False,
            )

        box = AgentFileSandbox(workspace_id=workspace_id)
        try:
            local_pptx = box.resolve(rel)
        except AgentFileSandbox.Error as exc:
            return json.dumps(
                {"ok": False, "error": exc.message, "code": "validation_error"},
                ensure_ascii=False,
            )

        expected_list: list[str] | None = None
        if expected_texts.strip():
            try:
                parsed = json.loads(expected_texts)
                if isinstance(parsed, list):
                    expected_list = [str(item) for item in parsed]
                else:
                    expected_list = [str(parsed)]
            except json.JSONDecodeError:
                expected_list = [line.strip() for line in expected_texts.splitlines() if line.strip()]

        result = await asyncio.to_thread(
            validate_ppt_output,
            local_pptx,
            expected_slide_count=expected_slide_count,
            expected_texts=expected_list,
        )
        return json.dumps({"ok": True, **result}, ensure_ascii=False)

    @tool
    async def draft_ppt_outline(
        brief: str,
        slide_count: int | None = None,
        output_path: str = "outlines/draft.json",
        source_md_path: str = "",
        project_dir: str = "ppt/default",
    ) -> str:
        """根据自然语言 brief 生成结构化 PPT 大纲 JSON 并写入沙箱。"""

        del project_dir  # reserved for future default output_path wiring

        if not brief.strip():
            return json.dumps(
                {"ok": False, "error": "brief is required", "code": "validation_error"},
                ensure_ascii=False,
            )
        if chat_model is None:
            return json.dumps(
                {"ok": False, "error": "chat model unavailable", "code": "model_unavailable"},
                ensure_ascii=False,
            )

        user_parts = [f"Brief:\n{brief.strip()}"]
        if slide_count is not None and slide_count > 0:
            user_parts.append(f"期望页数约 {slide_count} 页（含封面则 meta+slides）。")

        source_rel = (source_md_path or "").strip()
        if source_rel:
            box = AgentFileSandbox(workspace_id=workspace_id)
            try:
                read_result = await box.read_file_async(source_rel)
            except AgentFileSandbox.Error as exc:
                return json.dumps(
                    {"ok": False, "error": exc.message, "code": exc.code},
                    ensure_ascii=False,
                )
            if not read_result.get("ok"):
                return json.dumps(read_result, ensure_ascii=False)
            source_text = str(read_result.get("content", "")).strip()
            if source_text:
                user_parts.append(f"Source material:\n{source_text}")

        try:
            response = await chat_model.ainvoke(
                [
                    SystemMessage(content=_OUTLINE_SYSTEM),
                    HumanMessage(content="\n\n".join(user_parts)),
                ]
            )
        except Exception as exc:
            return json.dumps(
                {"ok": False, "error": str(exc), "code": "model_unavailable"},
                ensure_ascii=False,
            )

        text = response.content
        if isinstance(text, list):
            text = "\n".join(
                p.get("text", "") if isinstance(p, dict) else str(p) for p in text
            )
        if not isinstance(text, str):
            return json.dumps(
                {"ok": False, "error": "empty model response", "code": "invalid_json"},
                ensure_ascii=False,
            )

        try:
            raw = extract_json_object(text)
        except (json.JSONDecodeError, TypeError) as exc:
            return json.dumps(
                {"ok": False, "error": str(exc), "code": "invalid_json"},
                ensure_ascii=False,
            )

        try:
            doc = validate_outline_dict(raw)
        except Exception as exc:
            return json.dumps(
                {"ok": False, "error": str(exc), "code": "validation_error"},
                ensure_ascii=False,
            )

        outline = doc.model_dump()
        try:
            write_result = await _sandbox_write_json(workspace_id, output_path, outline)
        except AgentFileSandbox.Error as exc:
            return json.dumps(
                {"ok": False, "error": exc.message, "code": "write_failed"},
                ensure_ascii=False,
            )
        if not write_result.get("ok"):
            return json.dumps(write_result, ensure_ascii=False)

        preview = []
        for slide in outline.get("slides", []):
            if not isinstance(slide, dict):
                continue
            preview.append(
                {
                    "pageTitle": slide.get("pageTitle", slide.get("title", "")),
                    "itemCount": len(slide.get("items", []) or []),
                    "hasImages": bool(slide.get("images")) or bool(slide.get("hasImage")),
                }
            )

        return json.dumps(
            {
                "ok": True,
                "path": str(write_result.get("path", output_path)),
                "slides_count": len(outline.get("slides", [])),
                "preview": preview,
            },
            ensure_ascii=False,
        )

    @tool
    async def generate_ppt(
        outline: str = "",
        outline_path: str = "",
        output_path: str = "output/presentation.pptx",
        layout_mode: str = "hybrid",
        engine: str = "layout_fill",
        include_notes: bool = True,
        transition: str = "fade",
        template_path: str = "",
        fill_plan_path: str = "",
        project_dir: str = "ppt/default",
        design_spec_path: str = "",
        spec_lock_path: str = "",
        svg_dir: str = "",
    ) -> str:
        """从结构化大纲生成 pptx 文件到沙箱。outline 与 outline_path 二选一；svg_design 可无大纲。"""

        project = (project_dir or "ppt/default").strip().strip("/") or "ppt/default"
        engine_mode = (engine or "layout_fill").strip() or "layout_fill"

        mode = (layout_mode or "hybrid").strip().lower()
        if mode not in {"hybrid", "rule"}:
            mode = "hybrid"

        raw_outline: dict[str, Any] | None = None
        if outline_path.strip():
            box = AgentFileSandbox(workspace_id=workspace_id)
            try:
                read_result = await box.read_file_async(outline_path.strip())
            except AgentFileSandbox.Error as exc:
                return json.dumps(
                    {"ok": False, "error": exc.message, "code": "outline_invalid"},
                    ensure_ascii=False,
                )
            if not read_result.get("ok"):
                return json.dumps(read_result, ensure_ascii=False)
            try:
                raw_outline = json.loads(str(read_result.get("content", "")))
            except json.JSONDecodeError as exc:
                return json.dumps(
                    {"ok": False, "error": str(exc), "code": "outline_invalid"},
                    ensure_ascii=False,
                )
        elif outline.strip():
            try:
                raw_outline = json.loads(outline)
            except json.JSONDecodeError as exc:
                return json.dumps(
                    {"ok": False, "error": str(exc), "code": "outline_invalid"},
                    ensure_ascii=False,
                )
        elif engine_mode == "svg_design":
            raw_outline = {"slides": []}
        else:
            return json.dumps(
                {
                    "ok": False,
                    "error": "outline or outline_path is required",
                    "code": "outline_invalid",
                },
                ensure_ascii=False,
            )

        if not isinstance(raw_outline, dict):
            return json.dumps(
                {"ok": False, "error": "outline must be a JSON object", "code": "outline_invalid"},
                ensure_ascii=False,
            )

        if engine_mode != "svg_design":
            try:
                validate_outline_dict(raw_outline)
            except Exception as exc:
                return json.dumps(
                    {"ok": False, "error": str(exc), "code": "outline_invalid"},
                    ensure_ascii=False,
                )

        transition_mode = (transition or "fade").strip().lower()
        if transition_mode not in {"fade", "none", "keep"}:
            transition_mode = "fade"

        box = AgentFileSandbox(workspace_id=workspace_id)

        resolved_template: Path | None = None
        template_rel = (template_path or "").strip()
        if template_rel:
            try:
                resolved_template = box.resolve(template_rel)
            except AgentFileSandbox.Error as exc:
                return json.dumps(
                    {"ok": False, "error": exc.message, "code": "template_missing"},
                    ensure_ascii=False,
                )
        elif engine_mode == "template_fill":
            resolved_template = DEFAULT_TEMPLATE_PATH

        resolved_fill_plan: Path | None = None
        fill_plan_rel = (fill_plan_path or "").strip()
        if fill_plan_rel:
            try:
                resolved_fill_plan = box.resolve(fill_plan_rel)
            except AgentFileSandbox.Error as exc:
                return json.dumps(
                    {"ok": False, "error": exc.message, "code": "fill_plan_invalid"},
                    ensure_ascii=False,
                )

        resolved_svg_dir: Path | None = None
        svg_dir_rel = (svg_dir or "").strip()
        if svg_dir_rel:
            try:
                resolved_svg_dir = box.resolve(svg_dir_rel)
            except AgentFileSandbox.Error as exc:
                return json.dumps(
                    {"ok": False, "error": exc.message, "code": "svg_export_failed"},
                    ensure_ascii=False,
                )

        resolved_design_spec: Path | None = None
        design_spec_rel = (design_spec_path or "").strip()
        if design_spec_rel:
            try:
                resolved_design_spec = box.resolve(design_spec_rel)
            except AgentFileSandbox.Error as exc:
                return json.dumps(
                    {"ok": False, "error": exc.message, "code": "design_spec_missing"},
                    ensure_ascii=False,
                )

        resolved_spec_lock: Path | None = None
        spec_lock_rel = (spec_lock_path or "").strip()
        if spec_lock_rel:
            try:
                resolved_spec_lock = box.resolve(spec_lock_rel)
            except AgentFileSandbox.Error as exc:
                return json.dumps(
                    {"ok": False, "error": exc.message, "code": "design_spec_missing"},
                    ensure_ascii=False,
                )

        try:
            result = await generate_presentation(
                raw_outline,
                workspace_id=workspace_id,
                output_path=output_path,
                layout_mode=mode,
                chat_model=chat_model if mode == "hybrid" else None,
                engine=engine_mode,
                include_notes=include_notes,
                transition=transition_mode,
                template_path=resolved_template,
                fill_plan_path=resolved_fill_plan,
                project_dir=project,
                svg_dir=resolved_svg_dir,
                design_spec_path=resolved_design_spec,
                spec_lock_path=resolved_spec_lock,
            )
        except PptGenerateError as exc:
            return json.dumps(
                {"ok": False, "error": exc.message, "code": exc.code},
                ensure_ascii=False,
            )
        except AgentFileSandbox.Error as exc:
            return json.dumps(
                {"ok": False, "error": exc.message, "code": exc.code},
                ensure_ascii=False,
            )

        return json.dumps(result, ensure_ascii=False)

    validate_ppt_output_tool.__name__ = "validate_ppt_output"

    return [
        ingest_ppt_source,
        analyze_ppt_template,
        check_ppt_fill_plan,
        validate_ppt_output_tool,
        draft_ppt_outline,
        generate_ppt,
    ]
