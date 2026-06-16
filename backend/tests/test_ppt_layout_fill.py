"""Integration tests for layout_fill PPT generation."""

from __future__ import annotations

import json
import uuid
from pathlib import Path

import pytest
from pptx import Presentation

from app.agent.infrastructure.agent_file_sandbox import AgentFileSandbox
from app.agent.skills.ppt.pptmaker.generate import generate_presentation

FIXTURES = Path(__file__).resolve().parent / "fixtures" / "ppt"


@pytest.fixture
def agent_files_root(tmp_path, monkeypatch):
    """Point agent file storage at a temporary directory."""

    monkeypatch.setattr("app.config.resolve_agent_files_root", lambda: tmp_path)
    return tmp_path


@pytest.mark.asyncio
async def test_generate_presentation_rule_mode_creates_pptx(agent_files_root) -> None:
    """layout_mode=rule produces a non-empty pptx in sandbox."""

    workspace_id = uuid.uuid4()
    box = AgentFileSandbox(workspace_id=workspace_id)
    outline = json.loads((FIXTURES / "outline_minimal.json").read_text(encoding="utf-8"))

    result = await generate_presentation(
        outline,
        workspace_id=workspace_id,
        output_path="output/test.pptx",
        layout_mode="rule",
        chat_model=None,
    )
    assert result["ok"] is True
    dest = box.resolve(result["output_path"])
    assert dest.is_file()
    prs = Presentation(str(dest))
    assert len(prs.slides) >= 2
    assert result["pages"]
    assert result.get("engine") == "layout_fill"


@pytest.mark.asyncio
async def test_generate_presentation_clears_template_placeholder_text(agent_files_root) -> None:
    """Unfilled template hint text like 要点标题 is cleared from generated slides."""

    workspace_id = uuid.uuid4()
    box = AgentFileSandbox(workspace_id=workspace_id)
    outline = {
        "meta": {"title": "测试", "subtitle": "副标题"},
        "slides": [
            {
                "pageTitle": "目录",
                "items": ["章节一", "章节二", "章节三"],
            }
        ],
    }

    result = await generate_presentation(
        outline,
        workspace_id=workspace_id,
        output_path="output/toc_test.pptx",
        layout_mode="rule",
        chat_model=None,
    )
    assert result["ok"] is True
    dest = box.resolve(result["output_path"])
    prs = Presentation(str(dest))
    toc_slide = prs.slides[1]
    combined = "\n".join(
        shape.text_frame.text
        for shape in toc_slide.shapes
        if getattr(shape, "has_text_frame", False)
    )
    assert "要点标题" not in combined
    assert "要点正文" not in combined
    assert "章节一" in combined


@pytest.mark.asyncio
async def test_generate_presentation_writes_speaker_notes(agent_files_root) -> None:
    """speakerNotes in outline appear on the generated slide notes."""

    workspace_id = uuid.uuid4()
    box = AgentFileSandbox(workspace_id=workspace_id)
    outline = {
        "meta": {"title": "测试封面", "subtitle": "副标题"},
        "slides": [
            {
                "pageTitle": "备注页",
                "body": "正文内容",
                "speakerNotes": "这是演讲者备注文本",
            }
        ],
    }

    result = await generate_presentation(
        outline,
        workspace_id=workspace_id,
        output_path="output/notes_test.pptx",
        layout_mode="rule",
        chat_model=None,
        include_notes=True,
    )
    assert result["ok"] is True

    dest = box.resolve(result["output_path"])
    prs = Presentation(str(dest))
    content_slide = prs.slides[-1]
    notes_text = content_slide.notes_slide.notes_text_frame.text
    assert "这是演讲者备注文本" in notes_text
