"""Integration tests for template_fill engine."""

from __future__ import annotations

import json
import uuid
from pathlib import Path

import pytest
from pptx import Presentation

from app.agent.infrastructure.agent_file_sandbox import AgentFileSandbox
from app.agent.skills.ppt.pptmaker.constants import DEFAULT_TEMPLATE_PATH
from app.agent.skills.ppt.pptmaker.generate import generate_presentation
from app.agent.skills.ppt.template_fill.analyze import analyze_template
from app.agent.skills.ppt.template_fill.apply import apply_fill_plan
from app.agent.skills.ppt.template_fill.check_plan import check_fill_plan
from app.agent.skills.ppt.template_fill.plan_builder import outline_to_fill_plan
from app.agent.skills.ppt.tools import validate_ppt_output

FIXTURES = Path(__file__).resolve().parent / "fixtures" / "ppt"


@pytest.fixture
def agent_files_root(tmp_path, monkeypatch):
    """Point agent file storage at a temporary directory."""

    monkeypatch.setattr("app.config.resolve_agent_files_root", lambda: tmp_path)
    return tmp_path


def _create_mini_template(path: Path) -> dict:
    """Create a one-slide pptx with a known title shape and return slot metadata."""

    prs = Presentation()
    layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)
    title_shape = slide.shapes.title
    title_shape.text = "OLD TITLE"
    path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(path))
    return {"shape_id": title_shape.shape_id, "slide_index": 1}


def test_analyze_template_builtin_assets() -> None:
    """analyze_template returns non-empty slides with unique slot_ids."""

    library = analyze_template(DEFAULT_TEMPLATE_PATH)
    assert library["schema"] == "template_fill_pptx_library.v1"
    assert library["slides"]
    slot_ids = [slot["slot_id"] for slide in library["slides"] for slot in slide["slots"]]
    assert slot_ids
    assert len(slot_ids) == len(set(slot_ids))


def test_apply_fill_plan_replaces_title(tmp_path) -> None:
    """apply_fill_plan clones a slide and replaces title text."""

    template_path = tmp_path / "template_mini.pptx"
    meta = _create_mini_template(template_path)
    slot_id = f"s01_sh{meta['shape_id']}"
    plan = {
        "schema": "template_fill_pptx_plan.v1",
        "source_pptx": str(template_path),
        "slides": [
            {
                "source_slide": 1,
                "replacements": [{"slot_id": slot_id, "text": "NEW TITLE"}],
            }
        ],
    }
    output_path = tmp_path / "output.pptx"
    apply_fill_plan(template_path, plan, output_path)

    result = validate_ppt_output(output_path, expected_slide_count=1, expected_texts=["NEW TITLE"])
    assert result["passed"] is True
    assert "OLD TITLE" not in "\n".join(
        shape.text_frame.text
        for slide in Presentation(str(output_path)).slides
        for shape in slide.shapes
        if getattr(shape, "has_text_frame", False)
    )


def test_check_fill_plan_reports_overflow() -> None:
    """check_fill_plan emits overflow warnings for long replacement text."""

    library = {
        "slide_count": 1,
        "slides": [
            {
                "slide_index": 1,
                "slots": [
                    {
                        "slot_id": "s01_sh1",
                        "role": "title_candidate",
                        "geometry": {"w": 50.0},
                        "text_metrics": {"font_size_px": 24.0, "paragraph_count": 1},
                    }
                ],
            }
        ],
    }
    plan = {
        "slides": [
            {
                "source_slide": 1,
                "replacements": [{"slot_id": "s01_sh1", "text": "A" * 200}],
            }
        ],
    }
    report = check_fill_plan(library, plan)
    assert report["warnings"]


def test_outline_to_fill_plan_maps_cover_and_content() -> None:
    """outline_to_fill_plan builds replacements from outline fields."""

    library = analyze_template(DEFAULT_TEMPLATE_PATH)
    outline = json.loads((FIXTURES / "outline_minimal.json").read_text(encoding="utf-8"))
    plan = outline_to_fill_plan(outline, library, source_pptx="template.pptx")
    assert plan["schema"] == "template_fill_pptx_plan.v1"
    assert plan["slides"]
    assert any(r.get("text") for slide in plan["slides"] for r in slide.get("replacements", []))


@pytest.mark.asyncio
async def test_generate_presentation_template_fill_engine(agent_files_root, tmp_path) -> None:
    """generate_presentation with engine=template_fill writes replaced pptx."""

    workspace_id = uuid.uuid4()
    template_path = tmp_path / "template_mini.pptx"
    meta = _create_mini_template(template_path)
    slot_id = f"s01_sh{meta['shape_id']}"
    fill_plan = {
        "schema": "template_fill_pptx_plan.v1",
        "source_pptx": str(template_path),
        "slides": [
            {
                "source_slide": 1,
                "replacements": [{"slot_id": slot_id, "text": "NEW TITLE"}],
                "notes": "speaker note",
            }
        ],
    }
    fill_plan_path = tmp_path / "fill_plan_mini.json"
    fill_plan_path.write_text(json.dumps(fill_plan, ensure_ascii=False), encoding="utf-8")

    outline = {"slides": [{"pageTitle": "NEW TITLE", "body": "content"}]}
    result = await generate_presentation(
        outline,
        workspace_id=workspace_id,
        output_path="output/template_fill_test.pptx",
        engine="template_fill",
        template_path=template_path,
        fill_plan_path=fill_plan_path,
    )
    assert result["ok"] is True
    assert result["engine"] == "template_fill"

    box = AgentFileSandbox(workspace_id=workspace_id)
    dest = box.resolve(result["output_path"])
    validation = validate_ppt_output(dest, expected_slide_count=1, expected_texts=["NEW TITLE"])
    assert validation["passed"] is True
