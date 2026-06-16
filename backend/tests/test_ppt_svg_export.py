"""Tests for svg_design pipeline: finalize and export."""

from __future__ import annotations

import uuid
from pathlib import Path

import pytest
from pptx import Presentation

from app.agent.infrastructure.agent_file_sandbox import AgentFileSandbox
from app.agent.skills.ppt.pptmaker.generate import generate_presentation
from app.agent.skills.ppt.svg_pipeline.export import export_svgs_to_pptx
from app.agent.skills.ppt.svg_pipeline.finalize import finalize_svg_pages
from app.agent.skills.ppt.tools import validate_ppt_output

FIXTURES = Path(__file__).resolve().parent / "fixtures" / "ppt"
SVG_FIXTURE = FIXTURES / "svg" / "page_01.svg"


@pytest.fixture
def agent_files_root(tmp_path, monkeypatch):
    """Point agent file storage at a temporary directory."""

    monkeypatch.setattr("app.config.resolve_agent_files_root", lambda: tmp_path)
    return tmp_path


def test_finalize_finds_page_01_svg() -> None:
    """finalize_svg_pages returns sorted page_*.svg paths."""

    svg_dir = SVG_FIXTURE.parent
    pages = finalize_svg_pages(svg_dir)
    assert pages
    assert pages[0].name == "page_01.svg"


def test_export_svgs_to_pptx_creates_editable_slide(tmp_path) -> None:
    """export_svgs_to_pptx writes a pptx with expected text content."""

    output_path = tmp_path / "svg_export.pptx"
    count = export_svgs_to_pptx([SVG_FIXTURE], output_path)
    assert count == 1
    assert output_path.is_file()

    validation = validate_ppt_output(
        output_path,
        expected_slide_count=1,
        expected_texts=["Hello SVG Slide"],
    )
    assert validation["passed"] is True

    prs = Presentation(str(output_path))
    assert len(prs.slides) == 1
    texts = [
        shape.text_frame.text
        for slide in prs.slides
        for shape in slide.shapes
        if getattr(shape, "has_text_frame", False) and shape.text_frame.text.strip()
    ]
    assert any("Hello SVG Slide" in text for text in texts)


@pytest.mark.asyncio
async def test_generate_presentation_svg_design_engine(agent_files_root, tmp_path) -> None:
    """generate_presentation with engine=svg_design exports sandbox SVG pages."""

    workspace_id = uuid.uuid4()
    box = AgentFileSandbox(workspace_id=workspace_id)
    svg_dir = box.resolve("ppt/default/svg")
    svg_dir.mkdir(parents=True, exist_ok=True)
    (svg_dir / "design_spec.md").write_text("# Design spec\n", encoding="utf-8")
    (svg_dir / "spec_lock.md").write_text("locked\n", encoding="utf-8")
    (svg_dir / "page_01.svg").write_text(SVG_FIXTURE.read_text(encoding="utf-8"), encoding="utf-8")

    result = await generate_presentation(
        {"slides": []},
        workspace_id=workspace_id,
        output_path="output/svg_design_test.pptx",
        engine="svg_design",
        project_dir="ppt/default",
    )
    assert result["ok"] is True
    assert result["engine"] == "svg_design"
    assert result["pages"] == [{"page": 1, "svg": "page_01.svg"}]

    dest = box.resolve(result["output_path"])
    validation = validate_ppt_output(dest, expected_slide_count=1, expected_texts=["Hello SVG Slide"])
    assert validation["passed"] is True
